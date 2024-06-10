# Necessary imports

import json
import os
import sys
import re
import boto3
import streamlit as st
from langchain_community.embeddings import BedrockEmbeddings
from langchain_community.llms.bedrock import Bedrock
import numpy as np
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_community.vectorstores import FAISS
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import io
from collections import defaultdict
import pandas as pd

# AWS Bedrock Clients
bedrock = boto3.client(service_name="bedrock-runtime", region_name="us-east-1")
bedrock_embeddings = BedrockEmbeddings(model_id="amazon.titan-embed-text-v1", client=bedrock)

# Data Ingestion
def data_ingestion():
    """
        Create a data folder in your python file directory and add there the PDFs to be searched from.
    """
    loader = PyPDFDirectoryLoader("data")
    documents = loader.load()
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    docs = text_splitter.split_documents(documents)
    return docs

# Vector Store Creation
def get_vector_store(docs):
    vectorstore_faiss = FAISS.from_documents(docs, bedrock_embeddings)
    vectorstore_faiss.save_local("faiss_index")

# Agents Definitions
class QueryUnderstandingAgent:
    def refine_query(self, query):
        # For now, simply return the query as-is. This can be expanded for more complex query processing.
        return query

# Retrieval Agent for the pipeline entry
class RetrievalAgent:
    def retrieve_documents(self, vectorstore, query, k=3):
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": k})
        docs = retriever.get_relevant_documents(query)
        return docs

# Multiple Documents in data folder processed together
class SummarizationAgent:
    def summarize_documents(self, docs):
        # Join all document texts to create a context for the LLM
        return " ".join([doc.page_content for doc in docs])

# Final output generation agent
class GenerationAgent:
    def __init__(self, llm, retriever):
        self.llm = llm
        self.retriever = retriever

    def generate_response(self, context, question):
        prompt_template = """
        
        <context>
        {context}
        </context>

        Question: {question}

        Assistant:
        """
        PROMPT = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
        
        qa_chain = RetrievalQA.from_chain_type(
            llm=self.llm,
            chain_type="stuff",
            retriever=self.retriever,
            return_source_documents=True,
            chain_type_kwargs={"prompt": PROMPT}
        )
        response = qa_chain({"context": context, "query": question})
        return response['result']

class VerificationAgent:
    def verify_response(self, response):
        # For simplicity, we assume the response is always correct. This can be expanded with more complex verification logic.
        return response

# LLM Setup Functions
def get_claude_llm():
    llm = Bedrock(model_id="ai21.j2-mid-v1", client=bedrock, model_kwargs={'maxTokens': 512})
    return llm

def get_llama2_llm():
    llm = Bedrock(model_id="meta.llama2-70b-chat-v1", client=bedrock, model_kwargs={'max_gen_len': 512})
    return llm

# Create PowerPoint Presentation
def create_presentation(query, response, visuals):
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Document Response"
    content = slide.placeholders[1]
    content.text = query

    # Content slide with response
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Answer"
    content = slide.placeholders[1]
    content.text = response

    # Add slides for visuals
    for visual in visuals:
        slide_layout = prs.slide_layouts[5]  # Using a blank layout for visuals
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Visual Representation"
        
        # Add the visual image
        image_stream = io.BytesIO()
        visual.savefig(image_stream, format='png')
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1), width=Inches(8))

    prs.save("ppt/response_presentation.pptx")

# Extract numerical data and their contexts from the response using Regex
def extract_data_with_context(response):
    sentences = response.split('.')
    patterns = {
        'percentages': re.compile(r'(\d+(\.\d+)?%)'),
        'dates': re.compile(r'(\b\d{4}\b)'),
        'monetary': re.compile(r'(\$\d+(,\d{3})*(\.\d{2})?)'),
        'general_numbers': re.compile(r'\b\d+\b')
    }
    
    extracted_data = defaultdict(list)
    
    for sentence in sentences:
        for label, pattern in patterns.items():
            matches = pattern.findall(sentence)
            if matches:
                context = sentence.strip()
                for match in matches:
                    value = match[0] if isinstance(match, tuple) else match
                    extracted_data[label].append((value, context))

    return extracted_data

# Generate visuals based on extracted data with context
def generate_visuals(response):
    visuals = []

    # Extracted data with context
    data = extract_data_with_context(response)
    
    # Generate bar chart for general numbers if found
    if data['general_numbers']:
        numbers = [(int(num), context) for num, context in data['general_numbers']]
        labels, values = zip(*[(context, num) for num, context in numbers])

        # Create a bar chart
        plt.figure(figsize=(10, 5))
        plt.barh(labels, values, color='blue')
        plt.xlabel('Values')
        plt.title('Numerical Data Representation')
        visuals.append(plt.gcf())
    
    # Generate pie chart for percentages if found
    if data['percentages']:
        percentages = [(float(p.strip('%')), context) for p, context in data['percentages']]
        labels, values = zip(*[(context, p) for p, context in percentages])

        # Create a pie chart
        plt.figure(figsize=(10, 5))
        plt.pie(values, labels=labels, autopct='%1.1f%%')
        plt.title('Percentage Distribution')
        visuals.append(plt.gcf())

    # Generate line chart for dates if found
    if data['dates']:
        dates = sorted((int(date), context) for date, context in data['dates'])
        labels, values = zip(*[(context, date) for date, context in dates])

        # Create a line chart
        plt.figure(figsize=(10, 5))
        plt.plot(labels, values, marker='o', linestyle='-')
        plt.xlabel('Year')
        plt.ylabel('Values')
        plt.title('Trend Over Years')
        visuals.append(plt.gcf())

    # Generate bar chart for monetary values if found
    if data['monetary']:
        monetary_values = [(float(value.strip('$').replace(',', '')), context) for value, context in data['monetary']]
        labels, values = zip(*[(context, value) for value, context in monetary_values])

        # Create a bar chart
        plt.figure(figsize=(10, 5))
        plt.barh(labels, values, color='green')
        plt.xlabel('Monetary Values ($)')
        plt.title('Monetary Values Representation')
        visuals.append(plt.gcf())
    
    return visuals

# Main Function
def main():
    st.set_page_config(page_title="App Solutions", page_icon=":books:")
    st.header("App Enterprise Document Solution")

    user_question = st.text_input("Ask a Question from the PDF Files")

    with st.sidebar:
        st.title("Update Or Create Vector Store:")
        if st.button("Vectors Update"):
            with st.spinner("Processing..."):
                docs = data_ingestion()
                get_vector_store(docs)
                st.success("Done")

    if user_question:
        llm = None
        if st.button("1: Use Claude LLM"):
            llm = get_claude_llm()
        elif st.button("2: Use Llama2 LLM"):
            llm = get_llama2_llm()

        if llm is not None:
            with st.spinner("Processing..."):
                faiss_index = FAISS.load_local("faiss_index", bedrock_embeddings, allow_dangerous_deserialization=True)

                # Initialize agents
                query_agent = QueryUnderstandingAgent()
                retrieval_agent = RetrievalAgent()
                summarization_agent = SummarizationAgent()
                retriever = faiss_index.as_retriever(search_type="similarity", search_kwargs={"k": 3})
                generation_agent = GenerationAgent(llm, retriever)
                verification_agent = VerificationAgent()

                # Process the query
                refined_query = query_agent.refine_query(user_question)
                docs = retrieval_agent.retrieve_documents(faiss_index, refined_query)
                context = summarization_agent.summarize_documents(docs)
                response = generation_agent.generate_response(context, refined_query)
                verified_response = verification_agent.verify_response(response)
                
                st.write("Chatbot Response:")
                st.write(verified_response)

                # Generate visuals
                visuals = generate_visuals(verified_response)

                # Create PowerPoint presentation
                create_presentation(user_question, verified_response, visuals)
                
                st.success("Done")

if __name__ == "__main__":
    main()
